from pathlib import Path

from .base import Loader, Section


class TextLoader(Loader):
    def load(self, path):
        return [Section(path.read_text(encoding="utf-8"))]


class PdfLoader(Loader):
    def load(self, path):
        from pypdf import PdfReader
        from ..ocr.ocr_service import OCRService
        result = []
        for number, page in enumerate(PdfReader(path).pages, 1):
            text = (page.extract_text() or "").strip()
            if not text:
                text = OCRService().pdf_page(path, number - 1)
            result.append(Section(text=text, page=number))
        return result


class DocxLoader(Loader):
    def load(self, path):
        from docx import Document
        return [Section("\n".join(p.text for p in Document(path).paragraphs if p.text.strip()))]


class PptxLoader(Loader):
    def load(self, path):
        from pptx import Presentation
        sections = []
        for number, slide in enumerate(Presentation(path).slides, 1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            sections.append(Section(text, page=number, section_path=f"Slide {number}"))
        return sections


class ImageLoader(Loader):
    def load(self, path):
        from ..ocr.ocr_service import OCRService
        return [Section(OCRService().image(path), page=1)]


def loader_for(path: Path) -> Loader:
    loaders = {".txt": TextLoader, ".md": TextLoader, ".pdf": PdfLoader, ".docx": DocxLoader, ".pptx": PptxLoader,
               ".png": ImageLoader, ".jpg": ImageLoader, ".jpeg": ImageLoader, ".tif": ImageLoader, ".tiff": ImageLoader}
    try:
        return loaders[path.suffix.casefold()]()
    except KeyError as exc:
        raise ValueError(f"Unsupported document type: {path.suffix}") from exc
